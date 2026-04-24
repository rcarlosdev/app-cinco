from __future__ import annotations

import math
import re
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "deliverables" / "presentacion_ia_corporativa"
ASSETS_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "Presentacion_IA_Corporativa_Ejecutiva.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "navy": (12, 33, 58),
    "blue": (28, 92, 153),
    "teal": (25, 145, 138),
    "mint": (76, 191, 170),
    "amber": (233, 167, 38),
    "orange": (230, 108, 56),
    "ink": (27, 32, 43),
    "slate": (90, 103, 120),
    "light": (244, 247, 250),
    "white": (255, 255, 255),
    "line": (219, 226, 234),
    "success": (68, 179, 126),
    "warning": (227, 161, 38),
    "danger": (214, 85, 77),
}


def rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def ensure_dirs() -> None:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)


def load_domain_statuses() -> list[dict[str, str]]:
    registry = ROOT / "backend" / "apps" / "ia_dev" / "domains" / "registry"
    rows: list[dict[str, str]] = []
    for path in sorted(registry.glob("*.domain.yaml")):
        text = path.read_text(encoding="utf-8")
        get = lambda key: (
            re.search(rf"^{re.escape(key)}:\s*(.+)$", text, flags=re.MULTILINE).group(1).strip()
            if re.search(rf"^{re.escape(key)}:\s*(.+)$", text, flags=re.MULTILINE)
            else ""
        )
        rows.append(
            {
                "domain": get("dominio"),
                "name": get("nombre_dominio"),
                "status": get("estado_dominio"),
                "maturity": get("nivel_madurez"),
                "goal": get("objetivo_negocio"),
            }
        )
    return rows


def vertical_gradient(width: int, height: int, top: tuple[int, int, int], bottom: tuple[int, int, int]) -> Image.Image:
    img = Image.new("RGB", (width, height), top)
    draw = ImageDraw.Draw(img)
    for y in range(height):
        ratio = y / max(1, height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line((0, y, width, y), fill=color)
    return img


def create_cover_background() -> Path:
    width, height = 1600, 900
    img = vertical_gradient(width, height, COLORS["navy"], (8, 65, 94))
    draw = ImageDraw.Draw(img)

    for i in range(18):
        alpha = 18 + i * 3
        overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
        odraw = ImageDraw.Draw(overlay)
        x = 120 + i * 75
        odraw.ellipse((x, 90 + (i % 5) * 55, x + 280, 370 + (i % 4) * 35), fill=(35, 196, 182, alpha))
        overlay = overlay.filter(ImageFilter.GaussianBlur(20))
        img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")

    draw = ImageDraw.Draw(img)
    nodes = [
        (1120, 180),
        (1330, 250),
        (1230, 430),
        (1450, 470),
        (1100, 620),
        (1330, 690),
        (930, 420),
    ]
    for idx, (x1, y1) in enumerate(nodes):
        for x2, y2 in nodes[idx + 1 :]:
            if abs(x1 - x2) + abs(y1 - y2) < 420:
                draw.line((x1, y1, x2, y2), fill=(110, 222, 212), width=3)
    for x, y in nodes:
        draw.ellipse((x - 18, y - 18, x + 18, y + 18), fill=(255, 255, 255))
        draw.ellipse((x - 9, y - 9, x + 9, y + 9), fill=(233, 167, 38))

    path = ASSETS_DIR / "cover_background.png"
    img.save(path)
    return path


def create_silos_image() -> Path:
    width, height = 1200, 700
    img = Image.new("RGB", (width, height), COLORS["light"])
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((50, 50, 1150, 650), radius=36, fill=(255, 255, 255), outline=COLORS["line"], width=3)

    titles = [
        ("RRHH", COLORS["teal"]),
        ("Logistica", COLORS["amber"]),
        ("Transporte", COLORS["orange"]),
        ("Operaciones", COLORS["blue"]),
    ]
    x = 95
    for title, color in titles:
        draw.rounded_rectangle((x, 160, x + 220, 520), radius=28, fill=(250, 252, 255), outline=color, width=5)
        draw.rounded_rectangle((x + 30, 190, x + 190, 250), radius=16, fill=color)
        draw.line((x + 55, 310, x + 165, 310), fill=COLORS["slate"], width=10)
        draw.line((x + 55, 355, x + 175, 355), fill=COLORS["line"], width=10)
        draw.line((x + 55, 400, x + 150, 400), fill=COLORS["line"], width=10)
        draw.line((x + 55, 445, x + 130, 445), fill=COLORS["line"], width=10)
        draw.text((x + 60, 205), title, fill=(255, 255, 255))
        x += 250

    center = (600, 355)
    draw.ellipse((center[0] - 95, center[1] - 95, center[0] + 95, center[1] + 95), fill=COLORS["navy"])
    draw.ellipse((center[0] - 55, center[1] - 55, center[0] + 55, center[1] + 55), fill=COLORS["mint"])
    draw.ellipse((center[0] - 16, center[1] - 16, center[0] + 16, center[1] + 16), fill=(255, 255, 255))

    for target in [(205, 335), (455, 335), (705, 335), (955, 335)]:
        draw.line((target[0], target[1], center[0], center[1]), fill=(125, 140, 160), width=4)

    path = ASSETS_DIR / "silos_vs_intelligence.png"
    img.save(path)
    return path


def create_capability_map_image() -> Path:
    width, height = 1200, 700
    img = vertical_gradient(width, height, (239, 248, 249), (245, 247, 251))
    draw = ImageDraw.Draw(img)
    center = (600, 350)
    draw.ellipse((430, 180, 770, 520), fill=COLORS["navy"])
    draw.ellipse((500, 250, 700, 450), fill=COLORS["teal"])
    draw.text((535, 310), "Copiloto\nCorporativo", fill=(255, 255, 255), align="center")

    spokes = [
        ("Gestion\nHumana", (170, 170), COLORS["teal"]),
        ("SST", (170, 530), COLORS["mint"]),
        ("Logistica", (1030, 170), COLORS["amber"]),
        ("Transporte", (1030, 530), COLORS["orange"]),
        ("Operaciones", (600, 90), COLORS["blue"]),
        ("Gerencia /\nKPI", (600, 610), COLORS["slate"]),
    ]

    for label, (x, y), color in spokes:
        draw.rounded_rectangle((x - 110, y - 62, x + 110, y + 62), radius=24, fill=(255, 255, 255), outline=color, width=5)
        draw.line((center[0], center[1], x, y), fill=color, width=4)
        draw.text((x - 48, y - 24), label, fill=COLORS["ink"], align="center")

    path = ASSETS_DIR / "capability_map.png"
    img.save(path)
    return path


def create_logistics_flow_image() -> Path:
    width, height = 1200, 520
    img = Image.new("RGB", (width, height), (248, 249, 251))
    draw = ImageDraw.Draw(img)
    steps = [
        ("Ingreso", COLORS["blue"]),
        ("Inventario", COLORS["teal"]),
        ("Entrega", COLORS["amber"]),
        ("Consumo", COLORS["orange"]),
        ("Saldo y\ncontrol", COLORS["navy"]),
    ]
    box_w = 180
    gap = 40
    x = 70
    y = 160
    for idx, (label, color) in enumerate(steps):
        draw.rounded_rectangle((x, y, x + box_w, y + 120), radius=26, fill=(255, 255, 255), outline=color, width=5)
        draw.rounded_rectangle((x + 50, y + 22, x + 130, y + 62), radius=12, fill=color)
        draw.text((x + 48, y + 76), label, fill=COLORS["ink"], align="center")
        if idx < len(steps) - 1:
            draw.line((x + box_w, y + 60, x + box_w + gap, y + 60), fill=(125, 140, 160), width=8)
            arrow_x = x + box_w + gap - 12
            draw.polygon([(arrow_x, y + 60), (arrow_x - 16, y + 48), (arrow_x - 16, y + 72)], fill=(125, 140, 160))
        x += box_w + gap

    draw.rounded_rectangle((100, 340, 1100, 450), radius=24, fill=(230, 243, 240), outline=COLORS["mint"], width=3)
    draw.text((130, 368), "Control esperado: materiales, herramientas, equipos, dotaciones, consumos, saldos, devoluciones, responsables y trazabilidad por proyecto o cuadrilla.", fill=COLORS["ink"])

    path = ASSETS_DIR / "logistics_flow.png"
    img.save(path)
    return path


def create_transport_flow_image() -> Path:
    width, height = 1200, 520
    img = Image.new("RGB", (width, height), (247, 248, 251))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((60, 60, 1140, 460), radius=30, fill=(255, 255, 255), outline=COLORS["line"], width=3)

    milestones = [
        ("Vehiculo", (160, 260)),
        ("Asignacion", (360, 180)),
        ("Personal", (360, 340)),
        ("Ruta", (620, 180)),
        ("Salida", (620, 340)),
        ("Consumo", (870, 180)),
        ("Control", (1030, 340)),
    ]
    for label, (x, y) in milestones:
        draw.ellipse((x - 55, y - 55, x + 55, y + 55), fill=(235, 244, 252), outline=COLORS["blue"], width=5)
        draw.text((x - 34, y - 10), label, fill=COLORS["ink"])
    connections = [(0, 1), (0, 2), (1, 3), (2, 4), (3, 5), (4, 6), (5, 6)]
    for a, b in connections:
        draw.line((*milestones[a][1], *milestones[b][1]), fill=COLORS["orange"], width=5)

    path = ASSETS_DIR / "transport_flow.png"
    img.save(path)
    return path


def create_roadmap_image() -> Path:
    width, height = 1200, 520
    img = vertical_gradient(width, height, (246, 249, 252), (236, 242, 248))
    draw = ImageDraw.Draw(img)
    cards = [
        ("Fase 1", "Consolidar lo actual", "Empleados y ausentismo con adopcion gerencial.", COLORS["teal"]),
        ("Fase 2", "Expandir dominios", "Logistica, transporte, operaciones, horas extra y viaticos.", COLORS["amber"]),
        ("Fase 3", "Empresa conectada", "KPIs transversales, alertas, decisiones y seguimiento.", COLORS["navy"]),
    ]
    x = 70
    for idx, (phase, title, body, color) in enumerate(cards):
        draw.rounded_rectangle((x, 90, x + 320, 390), radius=28, fill=(255, 255, 255), outline=color, width=5)
        draw.rounded_rectangle((x + 28, 120, x + 140, 170), radius=14, fill=color)
        draw.text((x + 48, 132), phase, fill=(255, 255, 255))
        draw.text((x + 28, 205), title, fill=COLORS["ink"])
        draw.text((x + 28, 255), body, fill=COLORS["slate"])
        if idx < len(cards) - 1:
            draw.line((x + 320, 240, x + 360, 240), fill=(118, 132, 150), width=7)
            draw.polygon([(x + 360, 240), (x + 342, 228), (x + 342, 252)], fill=(118, 132, 150))
        x += 390

    path = ASSETS_DIR / "roadmap.png"
    img.save(path)
    return path


def create_domain_maturity_chart(rows: list[dict[str, str]]) -> Path:
    selected = []
    priority = {
        "empleados": 0,
        "ausentismo": 1,
        "transporte": 2,
        "facturacion": 3,
        "comisiones": 4,
        "horas_extras": 5,
        "viaticos": 6,
    }
    for row in rows:
        if row["domain"] in priority:
            selected.append(row)
    selected.sort(key=lambda r: priority.get(r["domain"], 99))

    status_score = {"planned": 1, "partial": 2, "active": 3}
    maturity_bonus = {"initial": 0.2, "growing": 0.5, "mature": 0.9}

    labels = [row["name"] for row in selected]
    values = [status_score.get(row["status"], 1) + maturity_bonus.get(row["maturity"], 0.0) for row in selected]
    colors = [
        "#44b37e" if row["status"] == "active" else "#e6a126" if row["status"] == "partial" else "#9aa7b7"
        for row in selected
    ]

    fig, ax = plt.subplots(figsize=(10.8, 5.6), dpi=150)
    ax.barh(labels, values, color=colors, edgecolor="none")
    ax.set_xlim(0, 4.2)
    ax.set_xlabel("Estado del dominio y nivel de madurez")
    ax.set_title("Madurez actual de dominios empresariales")
    ax.set_facecolor("#f8fafc")
    fig.patch.set_facecolor("#f8fafc")
    ax.grid(axis="x", linestyle="--", alpha=0.25)
    ax.spines[["top", "right", "left", "bottom"]].set_visible(False)
    for idx, row in enumerate(selected):
        ax.text(values[idx] + 0.05, idx, f"{row['status']} / {row['maturity']}", va="center", fontsize=9, color="#28323e")
    fig.tight_layout()

    path = ASSETS_DIR / "domain_maturity_chart.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_absence_chart() -> Path:
    labels = ["Operaciones", "Mantenimiento", "Logistica", "Gestion Humana", "Administrativo"]
    values = [38, 21, 15, 9, 6]
    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    colors = ["#1c5c99", "#19918a", "#4cbfaa", "#e9a726", "#e66c38"]
    ax.bar(labels, values, color=colors)
    ax.set_title("Ejemplo referencial: concentracion de ausentismo por area")
    ax.set_ylabel("Casos")
    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "absence_area_chart.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_inventory_chart() -> Path:
    categories = ["Materiales", "Herramientas", "Equipos", "Dotaciones"]
    entregado = [88, 74, 62, 91]
    saldo = [35, 22, 17, 29]
    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    x = range(len(categories))
    ax.bar([i - 0.16 for i in x], entregado, width=0.32, label="Entregado", color="#1c5c99")
    ax.bar([i + 0.16 for i in x], saldo, width=0.32, label="Saldo", color="#e9a726")
    ax.set_xticks(list(x), categories)
    ax.set_title("Ejemplo referencial: control de entregas y saldo disponible")
    ax.legend(frameon=False)
    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "inventory_chart.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_transport_chart() -> Path:
    days = ["Lun", "Mar", "Mie", "Jue", "Vie", "Sab"]
    salidas = [26, 29, 33, 31, 36, 18]
    consumo = [74, 79, 88, 83, 91, 44]
    fig, ax1 = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    ax1.plot(days, salidas, color="#1c5c99", marker="o", linewidth=3)
    ax1.set_ylabel("Salidas")
    ax1.set_title("Ejemplo referencial: salidas y consumo de flota")
    ax1.grid(axis="y", linestyle="--", alpha=0.25)
    ax2 = ax1.twinx()
    ax2.plot(days, consumo, color="#e66c38", marker="o", linewidth=3)
    ax2.set_ylabel("Consumo")
    ax1.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "transport_chart.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def add_full_bg(slide, color_name: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_name)


def add_picture(slide, path: Path, left: float, top: float, width: float | None = None, height: float | None = None):
    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), **kwargs)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 20,
    bold: bool = False,
    color: str = "ink",
    align=PP_ALIGN.LEFT,
    font_name: str = "Aptos",
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    *,
    font_size: int = 20,
    color: str = "ink",
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color)
        if idx:
            p.space_before = Pt(6)


def add_title(slide, title: str, subtitle: str | None = None, *, light: bool = False) -> None:
    title_color = "white" if light else "navy"
    subtitle_color = "light" if light else "slate"
    add_textbox(slide, 0.75, 0.45, 7.7, 0.6, title, font_size=28, bold=True, color=title_color)
    if subtitle:
        add_textbox(slide, 0.78, 1.03, 10.8, 0.38, subtitle, font_size=12, color=subtitle_color)


def add_footer(slide, text: str = "Material ejecutivo de referencia. Graficos ilustrativos salvo donde se indique lo contrario.") -> None:
    add_textbox(slide, 0.75, 7.03, 11.5, 0.22, text, font_size=9, color="slate")


def add_kpi_card(slide, left: float, top: float, width: float, height: float, title: str, value: str, accent: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("white")
    shape.line.color.rgb = rgb("line")
    shape.line.width = Pt(1.2)
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.18), Inches(height))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(accent)
    accent_bar.line.fill.background()
    add_textbox(slide, left + 0.32, top + 0.18, width - 0.45, 0.25, title, font_size=11, color="slate")
    add_textbox(slide, left + 0.32, top + 0.46, width - 0.45, 0.45, value, font_size=22, bold=True, color="navy")


def add_domain_cards(slide) -> None:
    cards = [
        ("Empleados", "Activo y creciendo", "Ya interpreta personal, supervisor, area, cargo, movil y tipo de labor.", "teal"),
        ("Ausentismo", "Activo y maduro", "Ya analiza vacaciones, incapacidades, licencias, permisos y concentraciones por estructura.", "mint"),
        ("Transporte", "Parcial y expandible", "Ya sirve como frente independiente de salidas y flota. Puede crecer a rutas, personal y consumo.", "amber"),
    ]
    x = 0.78
    for title, status, body, accent in cards:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.62), Inches(3.85), Inches(1.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb("white")
        shape.line.color.rgb = rgb("line")
        shape.line.width = Pt(1.1)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(1.84), Inches(1.6), Inches(0.32))
        pill.fill.solid()
        pill.fill.fore_color.rgb = rgb(accent)
        pill.line.fill.background()
        add_textbox(slide, x + 0.29, 1.89, 1.45, 0.2, status, font_size=9, bold=True, color="white")
        add_textbox(slide, x + 0.2, 2.26, 3.2, 0.34, title, font_size=18, bold=True, color="navy")
        add_textbox(slide, x + 0.2, 2.64, 3.3, 0.7, body, font_size=11, color="slate")
        x += 4.15


def build_presentation(asset_paths: dict[str, Path], domain_rows: list[dict[str, str]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture(slide, asset_paths["cover"], 0, 0, width=SLIDE_W, height=SLIDE_H)
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(0.76), Inches(1.8), Inches(0.32))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb("amber")
    accent.line.fill.background()
    add_textbox(slide, 0.88, 0.83, 1.35, 0.18, "Presentacion ejecutiva", font_size=10, bold=True, color="white")
    add_textbox(slide, 0.74, 1.44, 6.5, 1.25, "IA Corporativa Para Gestion,\nControl Y Decision", font_size=26, bold=True, color="white")
    add_textbox(slide, 0.78, 2.95, 5.7, 0.7, "Vision para escalar inteligencia empresarial por areas: gestion humana, SST, logistica, transporte, operaciones, control y gerencia.", font_size=16, color="light")
    add_textbox(slide, 0.8, 6.55, 5.3, 0.28, "Proyecto: app-cinco | Base actual validada en empleados y ausentismo", font_size=10, color="light")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "light")
    add_title(slide, "El problema actual", "La informacion existe, pero no siempre se convierte en accion")
    add_picture(slide, asset_paths["silos"], 6.9, 1.33, width=5.7)
    add_bullets(
        slide,
        0.82,
        1.62,
        5.4,
        4.2,
        [
            "La informacion suele estar repartida entre tablas, reportes, areas y personas.",
            "Muchos cruces siguen dependiendo de Excel, consultas manuales y conocimiento individual.",
            "La respuesta gerencial llega tarde cuando una novedad ya se volvio costo, retraso o riesgo.",
            "Los jefes de area necesitan ver su operacion completa, no datos aislados.",
        ],
        font_size=17,
    )
    add_kpi_card(slide, 0.86, 5.77, 1.6, 0.92, "Tiempo invertido", "Alto", "orange")
    add_kpi_card(slide, 2.62, 5.77, 1.85, 0.92, "Visibilidad cruzada", "Limitada", "amber")
    add_kpi_card(slide, 4.62, 5.77, 1.9, 0.92, "Velocidad de respuesta", "Variable", "blue")
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "La propuesta", "Un copiloto corporativo que entienda el lenguaje real de la empresa")
    add_picture(slide, asset_paths["capability_map"], 6.7, 1.35, width=5.8)
    add_bullets(
        slide,
        0.84,
        1.56,
        5.3,
        4.35,
        [
            "Permite consultar informacion empresarial en lenguaje natural y con criterio de negocio.",
            "Conecta personas, procesos, novedades, recursos, activos e indicadores.",
            "No reemplaza los sistemas actuales: los hace mas utilizables para decidir, controlar y actuar.",
            "Puede crecer por dominios, empezando por lo que mas duele y mas impacto genera.",
        ],
        font_size=17,
    )
    add_textbox(slide, 0.84, 5.8, 5.5, 0.48, "Resultado esperado: menos tiempo buscando datos y mas tiempo gestionando la operacion.", font_size=15, bold=True, color="navy")
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "light")
    add_title(slide, "Lo que ya esta demostrado", "Base real ya montada para crecer por procesos")
    add_domain_cards(slide)
    add_picture(slide, asset_paths["domain_chart"], 0.8, 3.82, width=6.15)
    add_textbox(slide, 7.25, 3.98, 5.0, 0.34, "Lo que esto significa para la gerencia", font_size=18, bold=True, color="navy")
    add_bullets(
        slide,
        7.28,
        4.36,
        4.7,
        2.2,
        [
            "Ya existe un punto de partida validado en gestion humana y ausentismo.",
            "El sistema ya entiende lenguaje interno como supervisor, area, carpeta, movil, vacaciones e incapacidad.",
            "La expansion a nuevos frentes no parte de cero; parte de una logica empresarial ya construida.",
        ],
        font_size=15,
    )
    add_footer(slide, "Madurez tomada del contexto actual del proyecto. Transporte aparece como dominio parcial; logistica de materiales se plantea como siguiente frente prioritario.")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Valor para Gestion Humana, RRHH y SST", "De consultar datos a gestionar personas con evidencia")
    add_picture(slide, asset_paths["absence_chart"], 7.02, 1.42, width=5.55)
    add_bullets(
        slide,
        0.82,
        1.58,
        5.55,
        4.35,
        [
            "Consultar personal activo, distribucion por area, cargo, supervisor o frente operativo.",
            "Detectar concentraciones de ausentismo y reincidencia antes de que afecten servicio y productividad.",
            "Ver vacaciones, incapacidades, permisos y licencias por fecha, lider o equipo.",
            "Priorizar acciones de SST con evidencia: focos de riesgo, recurrencia y poblacion critica.",
        ],
        font_size=16,
    )
    add_kpi_card(slide, 0.85, 5.98, 1.6, 0.86, "Visibilidad", "360°", "teal")
    add_kpi_card(slide, 2.6, 5.98, 1.9, 0.86, "Respuesta", "Mas rapida", "mint")
    add_kpi_card(slide, 4.62, 5.98, 2.05, 0.86, "Prevencion", "Mas temprana", "amber")
    add_footer(slide, "Grafico referencial para ilustrar el tipo de lectura ejecutiva que podria tener un jefe de area.")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "light")
    add_title(slide, "Valor para Logistica", "Materiales, herramientas, equipos y dotaciones bajo una sola lectura")
    add_picture(slide, asset_paths["logistics_flow"], 0.84, 1.56, width=6.2)
    add_picture(slide, asset_paths["inventory_chart"], 7.3, 1.72, width=5.0)
    add_bullets(
        slide,
        0.92,
        5.2,
        11.3,
        1.25,
        [
            "Permite controlar entrega, consumo, saldo, devoluciones, responsables, bodega, proyecto y cuadrilla.",
            "Ayuda a detectar faltantes, sobreconsumos, activos no retornados y diferencias entre inventario y realidad operativa.",
            "Impacta directamente costo, control interno, disponibilidad de recursos y productividad del frente de trabajo.",
        ],
        font_size=15,
    )
    add_footer(slide, "Logistica de materiales se presenta como siguiente dominio estrategico de expansion, separado de transporte.")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Valor para Transporte", "Control de flota, asignaciones, salidas, rutas y consumo")
    add_picture(slide, asset_paths["transport_flow"], 0.82, 1.6, width=6.25)
    add_picture(slide, asset_paths["transport_chart"], 7.1, 1.78, width=5.2)
    add_bullets(
        slide,
        0.88,
        5.18,
        11.2,
        1.35,
        [
            "Hace visible que vehiculo salio, con que personal, para que ruta, con que horario y con que consumo.",
            "Ayuda a medir uso de flota, tiempos muertos, desviaciones, costo por ruta y disciplina operativa.",
            "Conecta el control administrativo con la realidad del movimiento diario en campo.",
        ],
        font_size=15,
    )
    add_footer(slide, "Transporte ya existe como frente independiente en el proyecto y puede evolucionar a una operacion de flota mucho mas completa.")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "light")
    add_title(slide, "Valor para Operaciones y Productividad", "La mayor ganancia aparece cuando las areas se conectan")
    add_textbox(slide, 0.88, 1.55, 5.7, 0.45, "Escenario operativo real", font_size=19, bold=True, color="navy")
    add_bullets(
        slide,
        0.9,
        2.02,
        5.7,
        3.35,
        [
            "Una cuadrilla puede quedar impactada por ausentismo, falta de material o indisponibilidad de vehiculo.",
            "El jefe de operaciones necesita entender esas causas juntas, no por separado.",
            "El sistema puede priorizar frentes en riesgo, cuellos de botella y acciones correctivas inmediatas.",
        ],
        font_size=17,
    )
    # Process chain graphic
    boxes = [
        ("Personal", 7.0, 1.85, "teal"),
        ("Materiales", 9.2, 1.85, "amber"),
        ("Transporte", 7.0, 3.35, "orange"),
        ("Operacion", 9.2, 3.35, "blue"),
    ]
    for label, lx, ly, accent in boxes:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(ly), Inches(1.8), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb("white")
        shape.line.color.rgb = rgb(accent)
        shape.line.width = Pt(2.0)
        add_textbox(slide, lx + 0.18, ly + 0.26, 1.4, 0.25, label, font_size=16, bold=True, color="navy", align=PP_ALIGN.CENTER)
    for start, end in [((8.8, 2.3), (9.2, 2.3)), ((7.9, 2.75), (7.9, 3.35)), ((10.1, 2.75), (10.1, 3.35)), ((8.8, 3.8), (9.2, 3.8))]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1]))
        line.line.color.rgb = rgb("slate")
        line.line.width = Pt(2.5)
    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.18), Inches(2.72), Inches(1.55), Inches(1.1))
    center.fill.solid()
    center.fill.fore_color.rgb = rgb("navy")
    center.line.fill.background()
    add_textbox(slide, 8.36, 3.03, 1.18, 0.22, "Decision", font_size=15, bold=True, color="white", align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.0, 5.2, 5.3, 0.8, "Mensaje clave: la productividad sube cuando la empresa puede leer juntas las causas humanas, logisticas y operativas.", font_size=16, color="ink")
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Valor para Gerencia, Finanzas, Juridica e Indicadores", "Una sola capa de lectura ejecutiva sobre varios procesos")
    add_kpi_card(slide, 0.84, 1.72, 2.45, 1.05, "Costos y desviaciones", "Visibles", "orange")
    add_kpi_card(slide, 3.48, 1.72, 2.45, 1.05, "Cumplimiento", "Trazable", "blue")
    add_kpi_card(slide, 6.12, 1.72, 2.45, 1.05, "Riesgos", "Priorizables", "teal")
    add_kpi_card(slide, 8.76, 1.72, 2.75, 1.05, "Indicadores", "Conversacionales", "amber")
    add_bullets(
        slide,
        0.88,
        3.15,
        5.5,
        2.8,
        [
            "Finanzas puede leer consumos, variaciones y frentes de mayor impacto.",
            "Juridica puede apoyarse en trazabilidad, responsables, tiempos y evidencia operativa.",
            "Gerencia puede pedir KPI, desviaciones y causas probables sin esperar varios reportes separados.",
        ],
        font_size=16,
    )
    # Dashboard panel
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(3.02), Inches(5.35), Inches(2.75))
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb("light")
    panel.line.color.rgb = rgb("line")
    add_textbox(slide, 6.8, 3.22, 4.6, 0.3, "Cockpit ejecutivo esperado", font_size=17, bold=True, color="navy")
    add_kpi_card(slide, 6.82, 3.65, 1.55, 0.82, "Ausentismo", "Alerta", "orange")
    add_kpi_card(slide, 8.52, 3.65, 1.55, 0.82, "Dotacion", "Pendiente", "amber")
    add_kpi_card(slide, 10.22, 3.65, 1.45, 0.82, "Flota", "Ok", "success")
    add_textbox(slide, 6.85, 4.72, 4.7, 0.62, "Ejemplos de preguntas ejecutivas:\n1. Que areas tienen mayor impacto operativo esta semana?\n2. Donde se concentra el costo o el riesgo?\n3. Que debe intervenir primero la gerencia?", font_size=13, color="ink")
    add_footer(slide)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "navy")
    add_title(slide, "Siguiente paso", "De proyecto de IA a copiloto corporativo de empresa grande", light=True)
    add_picture(slide, asset_paths["roadmap"], 0.78, 1.48, width=7.1)
    add_bullets(
        slide,
        8.26,
        1.8,
        4.15,
        3.7,
        [
            "Fase 1: consolidar adopcion en empleados y ausentismo.",
            "Fase 2: abrir logistica de materiales y profundizar transporte.",
            "Fase 3: conectar operaciones, costos, indicadores y seguimiento directivo.",
            "Meta: que cada jefe de area tenga una capa inteligente de gestion sobre su proceso.",
        ],
        font_size=16,
        color="light",
    )
    add_textbox(slide, 8.28, 5.85, 4.0, 0.6, "Mensaje final", font_size=18, bold=True, color="white")
    add_textbox(slide, 8.28, 6.22, 4.2, 0.6, "No se trata de tener otro chat, sino de darle a la empresa una nueva forma de leer, controlar y dirigir su operacion.", font_size=14, color="light")
    add_footer(slide, "Presentacion generada a partir del contexto corporativo y de los dominios actualmente montados en el proyecto.")

    prs.save(PPTX_PATH)


def main() -> None:
    ensure_dirs()
    domain_rows = load_domain_statuses()
    asset_paths = {
        "cover": create_cover_background(),
        "silos": create_silos_image(),
        "capability_map": create_capability_map_image(),
        "logistics_flow": create_logistics_flow_image(),
        "transport_flow": create_transport_flow_image(),
        "roadmap": create_roadmap_image(),
        "domain_chart": create_domain_maturity_chart(domain_rows),
        "absence_chart": create_absence_chart(),
        "inventory_chart": create_inventory_chart(),
        "transport_chart": create_transport_chart(),
    }
    build_presentation(asset_paths=asset_paths, domain_rows=domain_rows)
    print(f"PPTX generado en: {PPTX_PATH}")


if __name__ == "__main__":
    main()
