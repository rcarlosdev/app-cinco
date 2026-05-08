from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BASE_ASSETS = ROOT / "deliverables" / "presentacion_ia_corporativa_v2" / "assets"

EXEC_DIR = ROOT / "deliverables" / "presentacion_documento_ejecutivo"
TECH_DIR = ROOT / "deliverables" / "presentacion_documento_tecnica"
EXEC_ASSETS = EXEC_DIR / "assets"
TECH_ASSETS = TECH_DIR / "assets"

EXEC_PPTX = EXEC_DIR / "Presentacion_Documento_Ejecutivo_Sistema_IA_Multiagente.pptx"
TECH_PPTX = TECH_DIR / "Presentacion_Guia_Tecnica_Runtime_IA_Multiagente.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "midnight": (13, 18, 38),
    "navy": (24, 30, 56),
    "violet": (114, 68, 226),
    "blue": (58, 112, 245),
    "teal": (28, 163, 151),
    "mint": (106, 214, 191),
    "amber": (234, 176, 55),
    "orange": (233, 111, 65),
    "ink": (35, 39, 51),
    "slate": (98, 107, 126),
    "line": (220, 227, 237),
    "soft": (245, 247, 251),
    "white": (255, 255, 255),
    "success": (72, 180, 126),
    "danger": (213, 86, 78),
}


def rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def ensure_dirs() -> None:
    EXEC_ASSETS.mkdir(parents=True, exist_ok=True)
    TECH_ASSETS.mkdir(parents=True, exist_ok=True)


def add_full_bg(slide, color_name: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color_name)
    shape.line.fill.background()


def add_bg_image(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def add_title(slide, title: str, subtitle: str | None = None, *, light: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = rgb("white" if light else "navy")
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.1), Inches(11.2), Inches(0.4))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = rgb("soft" if light else "slate")


def add_bullets(
    slide,
    bullets: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    color: str = "ink",
    level_1_size: int | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size if idx == 0 or level_1_size is None else level_1_size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(8)


def add_card(slide, *, left: float, top: float, width: float, height: float, title: str, body: str, accent: str) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("white")
    card.line.color.rgb = rgb("line")

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.12),
        Inches(height),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(accent)
    stripe.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.4), Inches(0.35))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = "Aptos Display"
    tr.font.bold = True
    tr.font.size = Pt(15)
    tr.font.color.rgb = rgb("navy")

    body_box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.6), Inches(width - 0.45), Inches(height - 0.8))
    bp = body_box.text_frame.paragraphs[0]
    br = bp.add_run()
    br.text = body
    br.font.name = "Aptos"
    br.font.size = Pt(11)
    br.font.color.rgb = rgb("ink")


def add_footer(slide, text: str, *, light: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = rgb("soft" if light else "slate")


def create_exec_value_chart() -> Path:
    labels = ["Tiempo de respuesta", "Dependencia tecnica", "Vision cruzada", "Accion inmediata"]
    antes = [86, 82, 34, 29]
    ahora = [32, 28, 83, 88]
    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    x = range(len(labels))
    ax.bar([i - 0.17 for i in x], antes, width=0.34, color="#bcc6d6", label="Antes")
    ax.bar([i + 0.17 for i in x], ahora, width=0.34, color="#6a4bed", label="Ahora")
    ax.set_xticks(list(x), labels)
    ax.set_ylim(0, 100)
    ax.set_ylabel("Indice referencial")
    ax.set_title("Cambio operativo esperado con IA multiagente")
    ax.legend(frameon=False)
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    out = EXEC_ASSETS / "valor_negocio_chart.png"
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def create_exec_delivery_chart() -> Path:
    labels = ["Datos", "Tablas", "Hallazgos", "Riesgos", "Recomendaciones", "Siguientes acciones"]
    values = [95, 78, 84, 69, 73, 81]
    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    ax.barh(labels, values, color=["#4f7cf0", "#20a397", "#6a4bed", "#ea6f41", "#eab037", "#6ad6bf"])
    ax.set_xlim(0, 100)
    ax.set_title("Tipo de salida que entrega el sistema")
    ax.set_xlabel("Cobertura funcional referencial")
    ax.grid(axis="x", linestyle="--", alpha=0.22)
    ax.spines[["top", "right", "left", "bottom"]].set_visible(False)
    fig.tight_layout()
    out = EXEC_ASSETS / "entregables_chart.png"
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def create_exec_flow_image() -> Path:
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), COLORS["soft"])
    draw = ImageDraw.Draw(img)
    boxes = [
        ("1. Pregunta", (80, 330, 320, 510), COLORS["blue"]),
        ("2. Entiende intencion", (400, 250, 700, 430), COLORS["violet"]),
        ("3. Busca fuente correcta", (780, 160, 1100, 340), COLORS["teal"]),
        ("4. Cruza datos", (780, 430, 1100, 610), COLORS["amber"]),
        ("5. Respuesta clara", (1180, 290, 1500, 470), COLORS["orange"]),
    ]
    for label, (x1, y1, x2, y2), color in boxes:
        draw.rounded_rectangle((x1, y1, x2, y2), radius=28, fill=(255, 255, 255), outline=color, width=6)
        draw.text((x1 + 26, y1 + 64), label, fill=COLORS["ink"])
    arrows = [
        ((320, 420), (400, 340)),
        ((700, 340), (780, 250)),
        ((700, 340), (780, 520)),
        ((1100, 250), (1180, 380)),
        ((1100, 520), (1180, 380)),
    ]
    for start, end in arrows:
        draw.line((*start, *end), fill=COLORS["slate"], width=8)
    out = EXEC_ASSETS / "flujo_negocio.png"
    img.save(out)
    return out


def create_tech_flow_image() -> Path:
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), (250, 251, 253))
    draw = ImageDraw.Draw(img)
    steps = [
        ("HTTP / CLI", 60, 350, 250, 470, "blue"),
        ("ChatApplicationService", 300, 330, 610, 490, "violet"),
        ("IntentArbitrationService", 660, 240, 980, 400, "teal"),
        ("SemanticBusinessResolver", 660, 500, 980, 660, "amber"),
        ("QueryExecutionPlanner", 1030, 330, 1330, 490, "orange"),
        ("sql_assisted | handler | fallback", 1370, 290, 1560, 530, "navy"),
    ]
    for label, x1, y1, x2, y2, color in steps:
        draw.rounded_rectangle((x1, y1, x2, y2), radius=26, fill=(255, 255, 255), outline=COLORS[color], width=6)
        draw.multiline_text((x1 + 20, y1 + 46), label, fill=COLORS["ink"], spacing=6)
    connectors = [
        ((250, 410), (300, 410)),
        ((610, 410), (660, 320)),
        ((610, 410), (660, 580)),
        ((980, 320), (1030, 410)),
        ((980, 580), (1030, 410)),
        ((1330, 410), (1370, 410)),
    ]
    for start, end in connectors:
        draw.line((*start, *end), fill=COLORS["slate"], width=8)
    out = TECH_ASSETS / "runtime_flow.png"
    img.save(out)
    return out


def create_source_of_truth_image() -> Path:
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), COLORS["soft"])
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((120, 170, 700, 710), radius=32, fill=(255, 255, 255), outline=COLORS["teal"], width=6)
    draw.rounded_rectangle((900, 170, 1480, 710), radius=32, fill=(255, 255, 255), outline=COLORS["amber"], width=6)
    draw.text((220, 220), "ai_dictionary", fill=COLORS["navy"])
    draw.multiline_text(
        (190, 300),
        "Fuente estructural unica\n\n- tablas reales\n- columnas reales\n- joins permitidos\n- sinonimos productivos\n- columnas autorizadas",
        fill=COLORS["ink"],
        spacing=12,
    )
    draw.text((1080, 220), "YAML de dominio", fill=COLORS["navy"])
    draw.multiline_text(
        (980, 300),
        "Narrativa y apoyo\n\n- contexto\n- vocabulario\n- ejemplos\n- reglas descriptivas\n\nNo debe definir estructura productiva",
        fill=COLORS["ink"],
        spacing=12,
    )
    draw.line((700, 440, 900, 440), fill=COLORS["slate"], width=8)
    draw.text((745, 395), "runtime", fill=COLORS["slate"])
    out = TECH_ASSETS / "source_of_truth.png"
    img.save(out)
    return out


def create_onboarding_chart() -> Path:
    steps = ["Registrar", "Auditar", "Preguntas reales", "Diagnose", "Verificar flujo"]
    values = [1, 1, 1, 1, 1]
    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    ax.bar(steps, values, color=["#4f7cf0", "#6a4bed", "#20a397", "#eab037", "#ea6f41"])
    ax.set_ylim(0, 1.3)
    ax.set_title("Secuencia correcta para onboardear un dominio")
    ax.set_ylabel("Paso obligatorio")
    ax.grid(axis="y", linestyle="--", alpha=0.20)
    ax.spines[["top", "right"]].set_visible(False)
    for i, label in enumerate(["dd_*", "audit", "negocio", "runtime", "sql_assisted"]):
        ax.text(i, 1.05, label, ha="center", va="bottom", fontsize=9, color="#30364a")
    fig.tight_layout()
    out = TECH_ASSETS / "onboarding_steps_chart.png"
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def create_decision_chart() -> Path:
    labels = ["sql_assisted", "handler", "legacy"]
    values = [9, 1, 0]
    fig, ax = plt.subplots(figsize=(7.8, 4.6), dpi=150)
    ax.bar(labels, values, color=["#6a4bed", "#20a397", "#bcc6d6"])
    ax.set_title("Estado validado del piloto cubierto")
    ax.set_ylabel("Casos")
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    out = TECH_ASSETS / "pilot_runtime_mix.png"
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def create_debug_chart() -> Path:
    labels = ["legacy_count", "runtime_only_fallback", "errores_sql", "insight_poor"]
    values = [0, 0, 0, 0]
    fig, ax = plt.subplots(figsize=(7.8, 4.6), dpi=150)
    ax.bar(labels, values, color="#20a397")
    ax.set_ylim(0, 1)
    ax.set_title("Indicadores a vigilar en pilot health/report")
    ax.set_ylabel("Valor esperado")
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    out = TECH_ASSETS / "debug_health_chart.png"
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return out


def build_executive_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    cover_bg = BASE_ASSETS / "cover_background_v2.png"
    domain_chart = BASE_ASSETS / "domain_maturity_chart_v2.png"
    absence_chart = BASE_ASSETS / "absence_area_chart_v2.png"
    transport_chart = BASE_ASSETS / "transport_chart_v2.png"
    purchase_chart = BASE_ASSETS / "purchase_impact_chart.png"
    logo = BASE_ASSETS / "corp_isotipo.png"

    flow_img = create_exec_flow_image()
    value_chart = create_exec_value_chart()
    delivery_chart = create_exec_delivery_chart()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(slide, cover_bg)
    slide.shapes.add_picture(str(logo), Inches(0.7), Inches(0.55), height=Inches(0.55))
    add_title(slide, "Sistema de IA Multiagente para Analisis Empresarial", "Version ejecutiva para liderazgo de negocio", light=True)
    hero = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.2), Inches(2.5))
    p = hero.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Preguntas directas.\nDatos reales.\nRespuestas accionables."
    r.font.name = "Aptos Display"
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = rgb("white")
    add_footer(slide, "IA multiagente | Presentacion ejecutiva", light=True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Resumen Ejecutivo", "Que hace el sistema hoy")
    add_bullets(
        slide,
        [
            "Permite hacer preguntas de negocio en lenguaje natural y responder con datos reales de la empresa.",
            "Reduce la dependencia de consultas manuales, cruces en Excel y apoyo tecnico permanente.",
            "Cruza informacion entre fuentes como empleados y ausentismo para devolver una lectura mas completa.",
            "Entrega no solo texto: tambien tablas, hallazgos, riesgos y recomendaciones.",
            "Ya esta validado en operacion controlada para ausentismo y empleados, con expansion prevista a transporte.",
        ],
        left=0.8,
        top=1.6,
        width=5.5,
        height=4.8,
        font_size=17,
    )
    slide.shapes.add_picture(str(value_chart), Inches(6.8), Inches(1.5), width=Inches(5.7))
    add_footer(slide, "El cambio clave es pasar de busqueda manual a respuesta guiada por datos.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Problema Que Resuelve", "Antes vs ahora")
    add_card(slide, left=0.8, top=1.6, width=3.7, height=3.8, title="Antes", body="Consultas manuales, dependencia de personas clave, reportes lentos y poca capacidad para cruzar fuentes a tiempo.", accent="danger")
    add_card(slide, left=4.8, top=1.6, width=3.7, height=3.8, title="Ahora", body="Un lider pregunta directamente y el sistema responde con datos reales, resumen ejecutivo y contexto para decidir.", accent="success")
    add_card(slide, left=8.8, top=1.6, width=3.7, height=3.8, title="Impacto", body="Mas velocidad, menos friccion, menos dependencia tecnica y una lectura mucho mas completa del negocio.", accent="violet")
    slide.shapes.add_picture(str(purchase_chart), Inches(1.0), Inches(5.6), width=Inches(11.4))
    add_footer(slide, "La ganancia no es solo automatizacion: es capacidad real de decision.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Como Funciona", "Explicado como un analista inteligente automatico")
    slide.shapes.add_picture(str(flow_img), Inches(0.7), Inches(1.4), width=Inches(12))
    add_footer(slide, "El usuario no necesita saber tablas ni nombres tecnicos.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Que Entrega El Sistema", "No solo responde: analiza")
    add_bullets(
        slide,
        [
            "Datos reales de la empresa",
            "Resumen ejecutivo entendible",
            "Tablas cuando se necesita detalle",
            "Hallazgos y alertas tempranas",
            "Riesgos detectados",
            "Recomendaciones y siguientes acciones",
        ],
        left=0.9,
        top=1.7,
        width=4.7,
        height=4.6,
        font_size=18,
    )
    slide.shapes.add_picture(str(delivery_chart), Inches(6.0), Inches(1.5), width=Inches(6.0))
    add_footer(slide, "La respuesta se diseña para que sea util para gestionar, no solo para leer.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Ejemplos Reales", "Casos representativos del uso actual")
    add_card(slide, left=0.8, top=1.6, width=3.8, height=2.1, title="Ejemplo 1", body="Que areas tienen mas ausentismo?\n\nRespuesta: ranking por area, interpretacion del foco y sugerencia de revision.", accent="violet")
    add_card(slide, left=0.8, top=4.0, width=3.8, height=2.1, title="Ejemplo 2", body="Cuantos empleados activos hay por cargo?\n\nRespuesta: conteo agrupado y lectura de concentraciones.", accent="teal")
    slide.shapes.add_picture(str(absence_chart), Inches(5.2), Inches(1.5), width=Inches(3.4))
    slide.shapes.add_picture(str(transport_chart), Inches(8.8), Inches(1.5), width=Inches(3.4))
    tech_box = slide.shapes.add_textbox(Inches(5.1), Inches(5.4), Inches(7.2), Inches(0.7))
    tp = tech_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "Ejemplo 3: cruces entre ausentismo y estructura de personal para detectar cargos o jefes con mayor concentracion de casos."
    tr.font.name = "Aptos"
    tr.font.size = Pt(13)
    tr.font.color.rgb = rgb("ink")
    add_footer(slide, "Los graficos son referentes visuales del tipo de salida esperada.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Estado Actual", "Lo que ya esta validado")
    add_bullets(
        slide,
        [
            "Operacion controlada con datos reales en ausentismo y empleados.",
            "Cruces entre fuentes funcionando dentro del alcance cubierto.",
            "Monitoreo de salud operativa y control de fallback.",
            "Base estable para replicar el modelo en nuevos dominios.",
        ],
        left=0.9,
        top=1.7,
        width=4.7,
        height=4.3,
        font_size=18,
    )
    slide.shapes.add_picture(str(domain_chart), Inches(6.0), Inches(1.45), width=Inches(6.1))
    status_box = slide.shapes.add_textbox(Inches(6.1), Inches(5.8), Inches(5.8), Inches(0.55))
    sp = status_box.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Lectura actual: ausentismo y empleados son la base validada; transporte es el siguiente frente natural."
    sr.font.name = "Aptos"
    sr.font.size = Pt(12)
    sr.font.color.rgb = rgb("slate")
    add_footer(slide, "El sistema ya paso de idea a fase validada.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "midnight")
    add_title(slide, "Proximos Pasos", "Expansion responsable y por etapas", light=True)
    add_card(slide, left=0.8, top=1.7, width=3.6, height=3.7, title="1. Consolidar", body="Mantener calidad y monitoreo sobre los dominios ya validados.", accent="mint")
    add_card(slide, left=4.85, top=1.7, width=3.6, height=3.7, title="2. Expandir", body="Abrir transporte y otras areas solo despues de cargar metadata y validar preguntas reales.", accent="amber")
    add_card(slide, left=8.9, top=1.7, width=3.6, height=3.7, title="3. Escalar", body="Mejorar precision, cobertura y capacidad analitica sin perder control del runtime.", accent="orange")
    quote = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.2), Inches(0.7))
    qp = quote.text_frame.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qr = qp.add_run()
    qr.text = "La oportunidad ya no es probar si funciona. La oportunidad es decidir donde genera mas valor despues."
    qr.font.name = "Aptos Display"
    qr.font.size = Pt(19)
    qr.font.bold = True
    qr.font.color.rgb = rgb("white")
    add_footer(slide, "Cierre ejecutivo", light=True)

    prs.save(EXEC_PPTX)


def build_technical_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    cover_bg = BASE_ASSETS / "cover_background_v2.png"
    logo = BASE_ASSETS / "corp_isotipo.png"

    flow_img = create_tech_flow_image()
    source_img = create_source_of_truth_image()
    onboarding_chart = create_onboarding_chart()
    decision_chart = create_decision_chart()
    debug_chart = create_debug_chart()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(slide, cover_bg)
    slide.shapes.add_picture(str(logo), Inches(0.7), Inches(0.55), height=Inches(0.55))
    add_title(slide, "Guia Tecnica De Onboarding Del Runtime IA Multiagente", "Version para developers y Staff Engineer review", light=True)
    hero = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(7.2), Inches(2.8))
    p = hero.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Menos heuristica manual.\nMas metadata correcta.\nUn solo runtime con autoridad."
    r.font.name = "Aptos Display"
    r.font.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = rgb("white")
    add_footer(slide, "Runtime IA multiagente | Onboarding tecnico", light=True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Flujo Unico Vigente", "El camino que no debemos romper")
    slide.shapes.add_picture(str(flow_img), Inches(0.75), Inches(1.45), width=Inches(12))
    footer_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.35))
    fp = footer_box.text_frame.paragraphs[0]
    fr = fp.add_run()
    fr.text = "HTTP/CLI -> ChatApplicationService -> IntentArbitrationService -> SemanticBusinessResolver -> QueryExecutionPlanner -> sql_assisted | handler moderno | runtime_fallback"
    fr.font.name = "Consolas"
    fr.font.size = Pt(11)
    fr.font.color.rgb = rgb("navy")
    add_footer(slide, "Toda extension nueva debe respetar esta ruta.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Que Hace Cada Componente", "Responsabilidades reales en produccion")
    add_card(slide, left=0.7, top=1.45, width=2.9, height=2.15, title="ChatApplicationService", body="Entry point real. Coordina contexto, memoria, arbitraje, planner, respuesta y observabilidad.", accent="violet")
    add_card(slide, left=3.85, top=1.45, width=2.9, height=2.15, title="IntentArbitrationService", body="GPT o fallback deterministico deciden la intencion y si aplica SQL, handler o fallback.", accent="teal")
    add_card(slide, left=7.0, top=1.45, width=2.9, height=2.15, title="SemanticBusinessResolver", body="Convierte intencion en contexto operativo usando ai_dictionary como autoridad estructural.", accent="amber")
    add_card(slide, left=10.15, top=1.45, width=2.45, height=2.15, title="QueryExecutionPlanner", body="Elige estrategia, valida SQL seguro y define la salida real.", accent="orange")
    add_card(slide, left=1.5, top=4.1, width=4.0, height=1.85, title="sql_assisted", body="Ruta preferida para analytics agregable con tablas, columnas y joins declarados.", accent="success")
    add_card(slide, left=5.8, top=4.1, width=3.2, height=1.85, title="handler", body="Solo para casos no expresables de forma razonable como SQL assisted.", accent="blue")
    add_card(slide, left=9.3, top=4.1, width=2.8, height=1.85, title="runtime_fallback", body="Contingencia encapsulada. No es plataforma de desarrollo.", accent="danger")
    add_footer(slide, "La autoridad moderna esta en application service + semantic resolver + planner.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Reglas De Arquitectura", "Lo que si y lo que no")
    add_card(slide, left=0.8, top=1.6, width=5.7, height=4.7, title="No hacer", body="- if palabra -> accion\n- logica estructural en YAML\n- agentes por tabla\n- SQL manual fuera del planner\n- rutas paralelas o reintroducir legacy como autoridad", accent="danger")
    add_card(slide, left=6.8, top=1.6, width=5.7, height=4.7, title="Si hacer", body="- ai_dictionary como source of truth\n- GPT arbitra la intencion\n- planner decide ejecucion\n- YAML solo narrativo\n- joins y sinonimos declarados antes de tocar codigo", accent="success")
    add_footer(slide, "Si una solucion nueva necesita mas hardcode, probablemente va en contra de la arquitectura.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Fuente Estructural Correcta", "ai_dictionary manda; YAML acompana")
    slide.shapes.add_picture(str(source_img), Inches(0.8), Inches(1.45), width=Inches(6.2))
    add_bullets(
        slide,
        [
            "Todo lo que el runtime necesita para compilar SQL debe existir en ai_dictionary.",
            "YAML puede describir contexto, ejemplos y vocabulario, pero no debe actuar como catalogo productivo.",
            "Si falta una columna o un join en ai_dictionary, el fix es metadata; no otro if.",
        ],
        left=7.3,
        top=1.8,
        width=4.7,
        height=3.6,
        font_size=17,
    )
    add_footer(slide, "La causa mas comun de mala extension es intentar compensar metadata incompleta con codigo manual.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Como Onboardear Un Dominio O Tabla", "Secuencia recomendada")
    slide.shapes.add_picture(str(onboarding_chart), Inches(0.9), Inches(1.55), width=Inches(5.9))
    add_bullets(
        slide,
        [
            "1. Registrar dominio, tabla, columnas, relaciones y sinonimos en dd_*.",
            "2. Ejecutar ia_dictionary_audit y corregir missing_* o yaml_structural_leaks.",
            "3. Probar preguntas reales del negocio, no prompts tecnicos.",
            "4. Ejecutar ia_runtime_diagnose con real-data.",
            "5. Verificar runtime_flow=sql_assisted y fallback_reason vacio cuando aplica.",
        ],
        left=6.9,
        top=1.75,
        width=5.1,
        height=4.5,
        font_size=16,
    )
    add_footer(slide, "Primero metadata, despues runtime. Nunca al reves.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Cuando Crear Handler", "Usarlo bien evita deuda")
    add_card(slide, left=0.9, top=1.75, width=5.2, height=3.0, title="Crear handler solo si", body="- no es realmente un query SQL\n- requiere logica especifica\n- necesita resolucion especial de entidad\n- el calculo no es razonable como SQL assisted", accent="amber")
    add_card(slide, left=6.35, top=1.75, width=5.2, height=3.0, title="No crear handler si", body="- es conteo, top, agrupacion o distribucion\n- se resuelve con joins declarados\n- el problema real es metadata faltante", accent="violet")
    slide.shapes.add_picture(str(decision_chart), Inches(3.0), Inches(5.05), width=Inches(7.2))
    add_footer(slide, "Si se puede resolver con SQL seguro, no abras otra ruta de negocio.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_title(slide, "Debugging Operativo", "Comandos y senales a revisar")
    add_bullets(
        slide,
        [
            "ia_runtime_diagnose: valida runtime_flow, compiler_used, fallback_reason, sql_assisted_count y legacy_count.",
            "ia_runtime_pilot_health: vigila status, legacy_count, runtime_only_fallback_count y errores_sql.",
            "ia_runtime_pilot_report: mira distribucion agregada de flow, fallback_reason y consultas problematicas.",
        ],
        left=0.9,
        top=1.75,
        width=5.7,
        height=4.2,
        font_size=17,
    )
    slide.shapes.add_picture(str(debug_chart), Inches(7.0), Inches(1.75), width=Inches(5.0))
    code_box = slide.shapes.add_textbox(Inches(7.0), Inches(5.2), Inches(5.0), Inches(1.0))
    tf = code_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Mirar siempre:"
    p1.font.name = "Aptos Display"
    p1.font.bold = True
    p1.font.size = Pt(13)
    p2 = tf.add_paragraph()
    p2.text = "runtime_flow | compiler_used | fallback_reason"
    p2.font.name = "Consolas"
    p2.font.size = Pt(12)
    p2.font.color.rgb = rgb("navy")
    add_footer(slide, "No tapes fallback_reason: es la pista de debugging mas valiosa.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_title(slide, "Ejemplo End To End", "Que areas tienen mas ausentismo?")
    add_card(slide, left=0.8, top=1.55, width=2.8, height=3.8, title="1. GPT arbitra", body="Detecta analytics_query sobre ausentismo y marca should_use_sql_assisted.", accent="violet")
    add_card(slide, left=3.95, top=1.55, width=2.8, height=3.8, title="2. Resolver semantica", body="SemanticBusinessResolver carga tabla de ausentismo, join con empleados y columna area.", accent="teal")
    add_card(slide, left=7.1, top=1.55, width=2.8, height=3.8, title="3. Planner", body="Genera SQL seguro con joins y agrupacion por area.", accent="amber")
    add_card(slide, left=10.25, top=1.55, width=2.2, height=3.8, title="4. Respuesta", body="runtime_flow=sql_assisted, datos reales, resumen y tabla.", accent="orange")
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(3.6),
        Inches(3.45),
        Inches(3.95),
        Inches(3.45),
    ).line.color.rgb = rgb("slate")
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(6.75),
        Inches(3.45),
        Inches(7.1),
        Inches(3.45),
    ).line.color.rgb = rgb("slate")
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(9.9),
        Inches(3.45),
        Inches(10.25),
        Inches(3.45),
    ).line.color.rgb = rgb("slate")
    add_footer(slide, "La salida correcta es razonar con metadata, no fijar una query manual por caso.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "midnight")
    add_title(slide, "Cierre Tecnico", "La arquitectura ya sabe razonar", light=True)
    add_bullets(
        slide,
        [
            "El developer no debe forzar la logica del negocio con mas heuristicas.",
            "La extension correcta alimenta ai_dictionary, valida runtime y deja que el planner haga su trabajo.",
            "Si una solucion nueva necesita fallback, YAML estructural o SQL manual, probablemente no esta alineada con la arquitectura actual.",
        ],
        left=0.95,
        top=1.8,
        width=10.8,
        height=3.6,
        font_size=19,
        color="white",
    )
    quote = slide.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.2), Inches(0.75))
    qp = quote.text_frame.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qr = qp.add_run()
    qr.text = "Metadata correcta primero. Runtime moderno despues. Hardcode nunca como estrategia."
    qr.font.name = "Aptos Display"
    qr.font.bold = True
    qr.font.size = Pt(20)
    qr.font.color.rgb = rgb("mint")
    add_footer(slide, "Cierre tecnico", light=True)

    prs.save(TECH_PPTX)


def main() -> None:
    ensure_dirs()
    build_executive_presentation()
    build_technical_presentation()
    print(f"Generated: {EXEC_PPTX}")
    print(f"Generated: {TECH_PPTX}")


if __name__ == "__main__":
    main()
