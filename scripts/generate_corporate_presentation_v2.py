from __future__ import annotations

import re
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from PIL import Image, ImageChops, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "deliverables" / "presentacion_ia_corporativa_v2"
ASSETS_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "Presentacion_IA_Corporativa_Ejecutiva_v2.pptx"

CORP_LOGO_PATH = Path(r"c:\Users\User\Claro drive\DEV\cinco_sas\2026\code_files\isotipo.png")
IA_ICON_PATH = Path(r"c:\Users\User\Downloads\logo_IA.jpeg")
IA_NAME_PATH = Path(r"c:\Users\User\Downloads\Nombre_IA.jpeg")

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "midnight": (13, 18, 38),
    "navy": (24, 30, 56),
    "violet": (114, 68, 226),
    "magenta": (217, 88, 235),
    "blue": (58, 112, 245),
    "teal": (28, 163, 151),
    "mint": (106, 214, 191),
    "amber": (234, 176, 55),
    "orange": (233, 111, 65),
    "ink": (35, 39, 51),
    "slate": (98, 107, 126),
    "line": (220, 227, 237),
    "soft": (245, 247, 251),
    "white": (255, 255, 255),
    "success": (72, 180, 126),
    "danger": (213, 86, 78),
}


def rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def ensure_dirs() -> None:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)


def load_domain_statuses() -> list[dict[str, str]]:
    registry = ROOT / "backend" / "apps" / "ia_dev" / "domains" / "registry"
    rows: list[dict[str, str]] = []
    for path in sorted(registry.glob("*.domain.yaml")):
        text = path.read_text(encoding="utf-8")

        def get(key: str) -> str:
            match = re.search(rf"^{re.escape(key)}:\s*(.+)$", text, flags=re.MULTILINE)
            return match.group(1).strip() if match else ""

        rows.append(
            {
                "domain": get("dominio"),
                "name": get("nombre_dominio"),
                "status": get("estado_dominio"),
                "maturity": get("nivel_madurez"),
                "goal": get("objetivo_negocio"),
            }
        )
    return rows


def vertical_gradient(width: int, height: int, top: tuple[int, int, int], bottom: tuple[int, int, int]) -> Image.Image:
    img = Image.new("RGB", (width, height), top)
    draw = ImageDraw.Draw(img)
    for y in range(height):
        ratio = y / max(1, height - 1)
        color = tuple(int(top[i] * (1 - ratio) + bottom[i] * ratio) for i in range(3))
        draw.line((0, y, width, y), fill=color)
    return img


def crop_to_visible(img: Image.Image) -> Image.Image:
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    alpha = img.getchannel("A")
    bbox = alpha.getbbox()
    if not bbox:
        return img
    return img.crop(bbox)


def knock_out_light_background(src: Path, dst: Path, threshold: int = 245) -> Path:
    img = Image.open(src).convert("RGBA")
    px = img.load()
    for y in range(img.height):
        for x in range(img.width):
            r, g, b, a = px[x, y]
            if r >= threshold and g >= threshold and b >= threshold:
                px[x, y] = (255, 255, 255, 0)
    crop_to_visible(img).save(dst)
    return dst


def knock_out_dark_background(src: Path, dst: Path, threshold: int = 22) -> Path:
    img = Image.open(src).convert("RGBA")
    px = img.load()
    for y in range(img.height):
        for x in range(img.width):
            r, g, b, a = px[x, y]
            if r <= threshold and g <= threshold and b <= threshold:
                px[x, y] = (0, 0, 0, 0)
    crop_to_visible(img).save(dst)
    return dst


def prepare_brand_assets() -> dict[str, Path]:
    corp_out = ASSETS_DIR / "corp_isotipo.png"
    ia_icon_out = ASSETS_DIR / "integr_ai_icon_transparent.png"
    ia_name_out = ASSETS_DIR / "integr_ai_name_transparent.png"

    corp_img = Image.open(CORP_LOGO_PATH).convert("RGBA")
    crop_to_visible(corp_img).save(corp_out)
    knock_out_dark_background(IA_ICON_PATH, ia_icon_out)
    knock_out_light_background(IA_NAME_PATH, ia_name_out)

    return {
        "corp_logo": corp_out,
        "ia_icon": ia_icon_out,
        "ia_name": ia_name_out,
    }


def create_cover_background() -> Path:
    width, height = 1600, 900
    img = vertical_gradient(width, height, COLORS["midnight"], (53, 31, 112))

    for idx in range(14):
        overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
        odraw = ImageDraw.Draw(overlay)
        x = 100 + idx * 95
        y = 70 + (idx % 4) * 85
        odraw.ellipse((x, y, x + 360, y + 260), fill=(120, 79, 234, 24 + idx * 4))
        overlay = overlay.filter(ImageFilter.GaussianBlur(18))
        img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")

    draw = ImageDraw.Draw(img)
    wave_points = [
        (920, 150),
        (1080, 210),
        (1210, 170),
        (1330, 280),
        (1180, 380),
        (1380, 465),
        (1070, 600),
        (1325, 705),
        (935, 420),
    ]
    for idx, (x1, y1) in enumerate(wave_points):
        for x2, y2 in wave_points[idx + 1 :]:
            if abs(x1 - x2) + abs(y1 - y2) < 430:
                draw.line((x1, y1, x2, y2), fill=(255, 255, 255), width=3)
    for x, y in wave_points:
        draw.ellipse((x - 16, y - 16, x + 16, y + 16), fill=(255, 255, 255))
        draw.ellipse((x - 7, y - 7, x + 7, y + 7), fill=(234, 176, 55))

    path = ASSETS_DIR / "cover_background_v2.png"
    img.save(path)
    return path


def create_domain_maturity_chart(rows: list[dict[str, str]]) -> Path:
    priority = {
        "empleados": 0,
        "ausentismo": 1,
        "transporte": 2,
        "facturacion": 3,
        "comisiones": 4,
        "horas_extras": 5,
        "viaticos": 6,
    }
    selected = [row for row in rows if row["domain"] in priority]
    selected.sort(key=lambda row: priority[row["domain"]])

    status_score = {"planned": 1.0, "partial": 2.0, "active": 3.0}
    maturity_bonus = {"initial": 0.2, "growing": 0.55, "mature": 0.95}
    labels = [row["name"] for row in selected]
    values = [status_score.get(row["status"], 1.0) + maturity_bonus.get(row["maturity"], 0.0) for row in selected]
    colors = []
    for row in selected:
        if row["status"] == "active":
            colors.append("#6a4bed")
        elif row["status"] == "partial":
            colors.append("#eab037")
        else:
            colors.append("#b4bfd0")

    fig, ax = plt.subplots(figsize=(10.6, 5.6), dpi=150)
    ax.barh(labels, values, color=colors, edgecolor="none")
    ax.set_xlim(0, 4.2)
    ax.set_xlabel("Estado y madurez del dominio")
    ax.set_title("Base actual y frentes de expansion")
    ax.grid(axis="x", linestyle="--", alpha=0.22)
    ax.spines[["top", "right", "left", "bottom"]].set_visible(False)
    ax.set_facecolor("#f8f9fc")
    fig.patch.set_facecolor("#f8f9fc")
    for idx, row in enumerate(selected):
        ax.text(values[idx] + 0.05, idx, f"{row['status']} / {row['maturity']}", va="center", fontsize=9, color="#30364a")
    fig.tight_layout()
    path = ASSETS_DIR / "domain_maturity_chart_v2.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_absence_chart() -> Path:
    labels = ["Operaciones", "Mantenimiento", "Logistica", "Gestion Humana", "Administrativo"]
    values = [38, 22, 16, 9, 5]
    fig, ax = plt.subplots(figsize=(8.6, 4.8), dpi=150)
    ax.bar(labels, values, color=["#6a4bed", "#4f7cf0", "#20a397", "#eab037", "#ea6f41"])
    ax.set_title("Ejemplo referencial: ausentismo por area")
    ax.set_ylabel("Casos")
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "absence_area_chart_v2.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_inventory_chart() -> Path:
    categories = ["Materiales", "Herramientas", "Equipos", "Dotaciones"]
    entregado = [89, 77, 64, 92]
    saldo = [34, 21, 16, 31]
    fig, ax = plt.subplots(figsize=(8.6, 4.8), dpi=150)
    x = range(len(categories))
    ax.bar([i - 0.16 for i in x], entregado, width=0.32, label="Entregado", color="#4f7cf0")
    ax.bar([i + 0.16 for i in x], saldo, width=0.32, label="Saldo", color="#eab037")
    ax.set_xticks(list(x), categories)
    ax.set_title("Ejemplo referencial: entregas y saldo disponible")
    ax.legend(frameon=False)
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "inventory_chart_v2.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_transport_chart() -> Path:
    days = ["Lun", "Mar", "Mie", "Jue", "Vie", "Sab"]
    salidas = [27, 30, 33, 31, 37, 19]
    consumo = [72, 78, 86, 82, 90, 43]
    fig, ax1 = plt.subplots(figsize=(8.6, 4.8), dpi=150)
    ax1.plot(days, salidas, marker="o", linewidth=3, color="#6a4bed")
    ax1.set_ylabel("Salidas")
    ax1.grid(axis="y", linestyle="--", alpha=0.22)
    ax2 = ax1.twinx()
    ax2.plot(days, consumo, marker="o", linewidth=3, color="#ea6f41")
    ax2.set_ylabel("Consumo")
    ax1.set_title("Ejemplo referencial: salidas y consumo de flota")
    ax1.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "transport_chart_v2.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_document_chart() -> Path:
    etapas = ["Lectura", "Extraccion", "Clasificacion", "Consolidacion", "Reporte"]
    manual = [8.0, 5.0, 4.0, 6.5, 5.0]
    asistido = [2.5, 1.7, 1.3, 2.2, 1.8]

    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    x = range(len(etapas))
    ax.bar([i - 0.17 for i in x], manual, width=0.34, color="#bcc6d6", label="Manual")
    ax.bar([i + 0.17 for i in x], asistido, width=0.34, color="#6a4bed", label="Con IA")
    ax.set_xticks(list(x), etapas)
    ax.set_ylabel("Horas por ciclo")
    ax.set_title("Ejemplo referencial: esfuerzo para producir un reporte")
    ax.legend(frameon=False)
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "document_automation_chart.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def create_purchase_impact_chart() -> Path:
    labels = ["Trazabilidad", "Velocidad", "Cumplimiento", "Control", "Visibilidad"]
    actual = [45, 52, 58, 49, 41]
    esperado = [88, 81, 86, 84, 92]

    fig, ax = plt.subplots(figsize=(8.8, 4.8), dpi=150)
    x = range(len(labels))
    ax.bar([i - 0.17 for i in x], actual, width=0.34, color="#bcc6d6", label="Actual")
    ax.bar([i + 0.17 for i in x], esperado, width=0.34, color="#20a397", label="Esperado")
    ax.set_xticks(list(x), labels)
    ax.set_ylim(0, 100)
    ax.set_ylabel("Indice referencial")
    ax.set_title("Ejemplo referencial: impacto esperado en compras")
    ax.legend(frameon=False)
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    path = ASSETS_DIR / "purchase_impact_chart.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def add_picture(slide, path: Path, left: float, top: float, width: float | None = None, height: float | None = None):
    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), **kwargs)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: str = "ink",
    align=PP_ALIGN.LEFT,
    font_name: str = "Aptos",
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    *,
    font_size: int = 16,
    color: str = "ink",
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {bullet}"
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color)
        if idx:
            p.space_before = Pt(6)


def add_full_bg(slide, color_name: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_name)


def add_title(slide, title: str, subtitle: str | None = None, *, light: bool = False) -> None:
    add_textbox(slide, 0.86, 0.46, 9.3, 0.52, title, font_size=27, bold=True, color="white" if light else "midnight")
    if subtitle:
        add_textbox(slide, 0.89, 1.0, 10.8, 0.3, subtitle, font_size=11, color="soft" if light else "slate")


def add_footer(slide, text: str) -> None:
    add_textbox(slide, 0.84, 7.0, 11.4, 0.2, text, font_size=8, color="slate")


def add_brand_marks(slide, assets: dict[str, Path], *, dark: bool = False) -> None:
    add_picture(slide, assets["corp_logo"], 0.34, 0.22, height=0.42)
    if dark:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.45), Inches(0.18), Inches(2.45), Inches(0.52))
        pill.fill.solid()
        pill.fill.fore_color.rgb = rgb("white")
        pill.line.fill.background()
        add_picture(slide, assets["ia_name"], 10.62, 0.26, height=0.28)
    else:
        add_picture(slide, assets["ia_name"], 10.35, 0.24, height=0.3)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    accent: str,
    title: str,
    body: str,
    fill: str = "white",
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb("line")
    shape.line.width = Pt(1.1)

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.18), Inches(top + 0.18), Inches(0.2), Inches(height - 0.36))
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb(accent)
    tag.line.fill.background()

    add_textbox(slide, left + 0.48, top + 0.2, width - 0.62, 0.28, title, font_size=15, bold=True, color="midnight")
    add_textbox(slide, left + 0.48, top + 0.58, width - 0.62, height - 0.7, body, font_size=11, color="slate")


def add_stat_card(slide, left: float, top: float, width: float, title: str, value: str, accent: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.92))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("white")
    shape.line.color.rgb = rgb("line")
    shape.line.width = Pt(1.0)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.16), Inches(0.92))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(accent)
    band.line.fill.background()

    add_textbox(slide, left + 0.26, top + 0.14, width - 0.35, 0.18, title, font_size=10, color="slate")
    add_textbox(slide, left + 0.26, top + 0.4, width - 0.35, 0.3, value, font_size=20, bold=True, color="midnight")


def add_stage_box(slide, left: float, top: float, width: float, title: str, accent: str, body: str | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("white")
    shape.line.color.rgb = rgb(accent)
    shape.line.width = Pt(2.0)
    add_textbox(slide, left + 0.12, top + 0.18, width - 0.24, 0.24, title, font_size=14, bold=True, color="midnight", align=PP_ALIGN.CENTER)
    if body:
        add_textbox(slide, left + 0.12, top + 0.5, width - 0.24, 0.32, body, font_size=9, color="slate", align=PP_ALIGN.CENTER)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color_name: str = "slate") -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color_name)
    line.line.width = Pt(2.6)


def add_file_tile(slide, left: float, top: float, code: str, accent: str, label: str) -> None:
    tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.18), Inches(0.96))
    tile.fill.solid()
    tile.fill.fore_color.rgb = rgb("white")
    tile.line.color.rgb = rgb("line")
    tile.line.width = Pt(1.0)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.18), Inches(top + 0.16), Inches(0.82), Inches(0.34))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(accent)
    badge.line.fill.background()
    add_textbox(slide, left + 0.27, top + 0.22, 0.65, 0.12, code, font_size=10, bold=True, color="white", align=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.58, 0.96, 0.18, label, font_size=9, color="slate", align=PP_ALIGN.CENTER)


def build_presentation(asset_paths: dict[str, Path], domain_rows: list[dict[str, str]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture(slide, asset_paths["cover_bg"], 0, 0, width=SLIDE_W, height=SLIDE_H)
    add_picture(slide, asset_paths["corp_logo"], 0.45, 0.42, height=0.62)

    white_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(1.24), Inches(5.55), Inches(4.62))
    white_card.fill.solid()
    white_card.fill.fore_color.rgb = rgb("white")
    white_card.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.55), Inches(1.55), Inches(0.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb("violet")
    accent.line.fill.background()
    add_textbox(slide, 1.14, 1.61, 1.2, 0.12, "Deck ejecutivo v2", font_size=9, bold=True, color="white")
    add_textbox(slide, 1.0, 2.0, 4.8, 0.82, "IntegrAI Corporativo\npara procesos empresariales", font_size=25, bold=True, color="midnight")
    add_textbox(
        slide,
        1.02,
        2.9,
        4.85,
        0.88,
        "Vision ejecutiva para gestion humana, SST, logistica, transporte, tratamiento de informacion, compras, operaciones y control gerencial.",
        font_size=14,
        color="slate",
    )
    add_picture(slide, asset_paths["ia_name"], 1.0, 4.02, height=0.48)
    add_textbox(slide, 1.02, 5.0, 4.7, 0.44, "Con base actual ya montada en empleados y ausentismo, y una ruta clara de expansion por dominio.", font_size=12, color="midnight")

    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.55), Inches(1.45), Inches(3.55), Inches(3.55))
    glow.fill.solid()
    glow.fill.fore_color.rgb = rgb("magenta")
    glow.line.fill.background()
    add_picture(slide, asset_paths["ia_icon"], 7.95, 1.28, height=4.1)
    add_textbox(slide, 8.52, 5.92, 3.75, 0.34, "Plataforma para leer, conectar y accionar la informacion de la empresa.", font_size=12, color="soft", align=PP_ALIGN.CENTER)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "El reto actual", "La informacion empresarial existe, pero muchas veces no esta lista para decidir")
    add_card(slide, 0.86, 1.58, 2.76, 1.45, accent="violet", title="Datos dispersos", body="Tablas, archivos, soportes, correos y reportes viven separados por area o sistema.")
    add_card(slide, 3.84, 1.58, 2.76, 1.45, accent="blue", title="Documentos sin lectura util", body="La empresa guarda informacion, pero muchas veces no la convierte en hallazgos ni seguimiento.")
    add_card(slide, 6.82, 1.58, 2.76, 1.45, accent="amber", title="Reportes manuales", body="Consolidar novedades, soportes y resultados toma tiempo y depende de pocas personas.")
    add_card(slide, 9.8, 1.58, 2.66, 1.45, accent="orange", title="Baja trazabilidad", body="Cuando aparece un atraso o sobrecosto, cuesta ver rapido la causa y el responsable.")

    add_textbox(slide, 0.9, 3.52, 4.8, 0.3, "Lo que esto genera", font_size=18, bold=True, color="midnight")
    add_bullets(
        slide,
        0.9,
        3.9,
        5.2,
        2.2,
        [
            "Mas tiempo buscando informacion y menos tiempo gestionando.",
            "Reportes que llegan cuando la novedad ya impacto la operacion.",
            "Dificultad para cruzar personas, recursos, soportes y procesos.",
        ],
    )

    # visual ladder
    steps = [
        ("Dato", 7.0, "violet"),
        ("Archivo", 8.25, "blue"),
        ("Proceso", 9.5, "amber"),
        ("Decision", 10.75, "teal"),
    ]
    for label, lx, accent_name in steps:
        add_stage_box(slide, lx, 3.95, 1.05, label, accent_name)
    add_arrow(slide, 8.05, 4.47, 8.25, 4.47)
    add_arrow(slide, 9.3, 4.47, 9.5, 4.47)
    add_arrow(slide, 10.55, 4.47, 10.75, 4.47)

    add_stat_card(slide, 0.9, 6.04, 1.65, "Velocidad", "Lenta", "danger")
    add_stat_card(slide, 2.72, 6.04, 1.9, "Visibilidad cruzada", "Parcial", "amber")
    add_stat_card(slide, 4.8, 6.04, 1.9, "Trazabilidad", "Variable", "blue")
    add_footer(slide, "Lectura ejecutiva del problema: no basta con almacenar informacion; hay que volverla visible, accionable y trazable.")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "La propuesta", "Un copiloto corporativo que entienda el lenguaje real del negocio")

    core = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.84), Inches(2.0), Inches(3.6), Inches(2.1))
    core.fill.solid()
    core.fill.fore_color.rgb = rgb("midnight")
    core.line.fill.background()
    add_textbox(slide, 5.35, 2.55, 2.6, 0.4, "IntegrAI\nCorporativo", font_size=22, bold=True, color="white", align=PP_ALIGN.CENTER)

    spokes = [
        ("Gestion Humana", 1.15, 1.55, "violet"),
        ("SST", 1.25, 4.45, "teal"),
        ("Logistica", 9.75, 1.55, "amber"),
        ("Transporte", 9.72, 4.45, "orange"),
        ("Compras", 4.95, 5.35, "blue"),
        ("Documentos e imagenes", 4.55, 0.95, "magenta"),
    ]
    for label, lx, ly, accent_name in spokes:
        bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(ly), Inches(2.25), Inches(0.86))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = rgb("soft")
        bubble.line.color.rgb = rgb(accent_name)
        bubble.line.width = Pt(2.0)
        add_textbox(slide, lx + 0.12, ly + 0.22, 2.0, 0.2, label, font_size=13, bold=True, color="midnight", align=PP_ALIGN.CENTER)
        add_arrow(slide, 6.63, 3.02, lx + 1.12, ly + 0.43, accent_name)

    add_bullets(
        slide,
        0.92,
        6.15,
        11.2,
        0.62,
        [
            "Consulta datos, interpreta archivos y soportes, y convierte hallazgos en seguimiento, control y accion.",
        ],
        font_size=14,
    )
    add_footer(slide, "La oportunidad no es tener otro chat: es tener una capa inteligente sobre varios procesos empresariales.")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Base actual ya construida", "Hay una base real para escalar por dominios y no empezar de cero")
    add_card(slide, 0.86, 1.56, 3.88, 1.7, accent="violet", title="Empleados", body="Ya interpreta personal, colaboradores, cedula, supervisor, area, cargo, carpeta, movil y tipo de labor.")
    add_card(slide, 4.92, 1.56, 3.88, 1.7, accent="teal", title="Ausentismo", body="Ya entiende vacaciones, incapacidades, licencias, permisos, reincidencia y lectura por periodo o estructura.")
    add_card(slide, 8.98, 1.56, 3.45, 1.7, accent="amber", title="Expansiones claras", body="Transporte, compras, logistica de materiales, horas extra, viaticos, comisiones, facturacion.")
    add_picture(slide, asset_paths["domain_chart"], 0.88, 3.55, width=6.15)
    add_bullets(
        slide,
        7.25,
        3.85,
        4.85,
        2.4,
        [
            "Ya existe una logica de negocio que se puede replicar por area.",
            "La empresa ya puede mostrar resultados reales en dos frentes de alto valor.",
            "La expansion puede hacerse por etapas, segun impacto y prioridad.",
        ],
        font_size=15,
    )
    add_footer(slide, "La lectura de madurez se apoya en los dominios empresariales actualmente definidos en el proyecto.")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Valor para Gestion Humana, RRHH y SST", "De consultar datos a gestionar personas con evidencia")
    add_picture(slide, asset_paths["absence_chart"], 7.05, 1.46, width=5.35)
    add_bullets(
        slide,
        0.9,
        1.62,
        5.7,
        4.3,
        [
            "Consultar personal activo por area, cargo, supervisor, carpeta o frente operativo.",
            "Detectar concentraciones de ausentismo y recurrencia antes de que afecten servicio y productividad.",
            "Ver vacaciones, incapacidades, permisos y licencias por fecha, lider, area o equipo.",
            "Dar a SST una base para priorizar focos de riesgo e intervenir con evidencia.",
        ],
        font_size=16,
    )
    add_stat_card(slide, 0.92, 6.0, 1.55, "Visibilidad", "360", "violet")
    add_stat_card(slide, 2.62, 6.0, 1.92, "Respuesta", "Mas rapida", "teal")
    add_stat_card(slide, 4.72, 6.0, 2.0, "Prevencion", "Mas temprana", "amber")
    add_footer(slide, "Grafico referencial para ilustrar el tipo de lectura ejecutiva que puede tener un jefe de area.")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Valor para Logistica", "Materiales, herramientas, equipos y dotaciones bajo una sola lectura")
    add_picture(slide, asset_paths["inventory_chart"], 7.15, 1.65, width=5.3)

    stages = [
        ("Ingreso", 0.92, "violet", "Recepcion y legalizacion"),
        ("Inventario", 2.48, "blue", "Existencia y ubicacion"),
        ("Entrega", 4.04, "amber", "Responsable y destino"),
        ("Consumo", 5.6, "orange", "Uso y reposicion"),
        ("Saldo / control", 7.16, "teal", "Trazabilidad y diferencia"),
    ]
    for title, lx, accent_name, body in stages:
        add_stage_box(slide, lx, 1.74, 1.34, title, accent_name, body)
        if lx < 7.0:
            add_arrow(slide, lx + 1.34, 2.26, lx + 1.52, 2.26)

    add_bullets(
        slide,
        0.95,
        3.45,
        5.95,
        2.5,
        [
            "Controla entrega, consumo, devoluciones, saldo, responsable, bodega, cuadrilla y proyecto.",
            "Ayuda a detectar faltantes, sobreconsumos, activos no retornados e inconsistencias entre soporte y realidad operativa.",
            "Impacta costo, disponibilidad, control interno y continuidad del frente de trabajo.",
        ],
        font_size=15,
    )
    add_footer(slide, "Aqui logistica se entiende como manejo de materiales, herramientas, equipos y dotaciones. Transporte se trata como dominio aparte.")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Valor para Transporte", "Control de vehiculos, asignaciones, salidas, rutas y consumo")
    add_picture(slide, asset_paths["transport_chart"], 7.1, 1.7, width=5.3)

    top_row = [
        ("Vehiculo", 0.95, 1.76, "violet"),
        ("Asignacion", 2.55, 1.76, "blue"),
        ("Ruta", 4.15, 1.76, "amber"),
        ("Consumo", 5.75, 1.76, "orange"),
    ]
    bottom_row = [
        ("Personal", 1.75, 3.02, "teal"),
        ("Salida", 3.55, 3.02, "magenta"),
        ("Control", 5.35, 3.02, "navy"),
    ]
    for label, lx, ly, accent_name in top_row + bottom_row:
        add_stage_box(slide, lx, ly, 1.26, label, accent_name)
    for start, end in [((2.21, 2.28), (2.55, 2.28)), ((3.81, 2.28), (4.15, 2.28)), ((5.41, 2.28), (5.75, 2.28)), ((3.18, 2.8), (2.38, 3.02)), ((4.78, 2.8), (4.18, 3.02)), ((6.38, 2.8), (5.98, 3.02))]:
        add_arrow(slide, start[0], start[1], end[0], end[1])

    add_bullets(
        slide,
        0.95,
        4.42,
        5.95,
        1.75,
        [
            "Permite ver que vehiculo salio, con que personal, para que ruta, con que horario y con que consumo.",
            "Ayuda a medir uso de flota, tiempos muertos, desvio operativo y costo por ruta o asignacion.",
        ],
        font_size=15,
    )
    add_footer(slide, "Transporte se presenta como frente independiente de control operativo de flota y asignaciones.")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Tratamiento de informacion", "Lectura de archivos, documentos e imagenes para volverlos accion empresarial")

    add_file_tile(slide, 0.94, 1.72, "PDF", "violet", "Contratos / actas")
    add_file_tile(slide, 2.25, 1.72, "XLS", "blue", "Bases / reportes")
    add_file_tile(slide, 3.56, 1.72, "IMG", "amber", "Fotos / soportes")
    add_file_tile(slide, 4.87, 1.72, "DOC", "teal", "Formatos / cartas")
    add_file_tile(slide, 6.18, 1.72, "MAIL", "magenta", "Correos / solicitudes")

    add_arrow(slide, 7.15, 2.18, 7.55, 2.18, "slate")
    add_stage_box(slide, 7.58, 1.62, 1.5, "Extraccion", "violet", "Fechas, valores, nombres, estados")
    add_arrow(slide, 9.08, 2.18, 9.48, 2.18, "slate")
    add_stage_box(slide, 9.52, 1.62, 1.5, "Clasificacion", "blue", "Tipo, area, prioridad, responsable")
    add_arrow(slide, 11.02, 2.18, 11.38, 2.18, "slate")
    add_stage_box(slide, 11.42, 1.62, 1.45, "Reporte", "teal", "Informe, alerta o seguimiento")

    add_bullets(
        slide,
        0.96,
        3.38,
        11.0,
        2.45,
        [
            "Puede leer soportes documentales y convertirlos en datos usables para control y gestion.",
            "Puede resumir documentos extensos, consolidar multiples archivos y detectar faltantes o inconsistencias.",
            "Permite que el conocimiento ya no se quede atrapado en carpetas, correos o imagenes sueltas.",
        ],
        font_size=15,
    )
    add_footer(slide, "Esta expansion complementa la base analitica actual y abre la puerta a una capa documental y multimodal.")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "De soportes a informes y reportes", "La lectura de documentos puede convertirse en seguimiento, control y reporteria ejecutiva")
    add_picture(slide, asset_paths["document_chart"], 6.9, 1.5, width=5.55)

    add_card(slide, 0.9, 1.65, 2.68, 1.28, accent="violet", title="RRHH", body="Incapacidades, permisos, hojas de vida, certificados, novedades y vencimientos.")
    add_card(slide, 0.9, 3.05, 2.68, 1.28, accent="teal", title="SST", body="Actas, investigaciones, inspecciones, hallazgos, responsables y fechas compromiso.")
    add_card(slide, 0.9, 4.45, 2.68, 1.28, accent="amber", title="Logistica", body="Remisiones, entregas, devoluciones, soportes de consumo y diferencias de inventario.")
    add_card(slide, 3.8, 1.65, 2.68, 1.28, accent="blue", title="Finanzas", body="Facturas, legalizaciones, viaticos, soportes de gasto y centro de costo.")
    add_card(slide, 3.8, 3.05, 2.68, 1.28, accent="orange", title="Juridica", body="Contratos, otrosi, oficios, requerimientos, fechas y obligaciones.")
    add_card(slide, 3.8, 4.45, 2.68, 1.28, accent="magenta", title="Gerencia", body="Consolidado de hallazgos, pendientes, alertas y resumen ejecutivo.")
    add_footer(slide, "Grafico referencial: ilustra la reduccion de esfuerzo para pasar de archivo disperso a reporte util.")

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Proceso de compras", "Ciclo actual que puede volverse totalmente visible y trazable")
    compra_stages = [
        ("Solicitud", 0.92, "violet", "Necesidad y requerimiento"),
        ("Aprobacion", 2.88, "blue", "Decision y visto bueno"),
        ("Cotizacion", 4.84, "amber", "Comparativo y proveedor"),
        ("Compra", 6.8, "orange", "Orden y ejecucion"),
        ("Estado de entrega", 8.76, "teal", "Seguimiento y recepcion"),
        ("Uso final", 10.72, "magenta", "Destino, consumo y cierre"),
    ]
    for idx, (title, lx, accent_name, body) in enumerate(compra_stages):
        add_stage_box(slide, lx, 2.05, 1.58, title, accent_name, body)
        if idx < len(compra_stages) - 1:
            add_arrow(slide, lx + 1.58, 2.57, lx + 1.83, 2.57)

    add_textbox(slide, 0.96, 4.1, 2.2, 0.22, "Aportes esperados en cada etapa", font_size=16, bold=True, color="midnight")
    add_card(slide, 0.95, 4.45, 1.85, 1.42, accent="violet", title="Solicitud", body="Lectura de formatos, correo o requerimiento y estructuracion automatica.")
    add_card(slide, 3.02, 4.45, 1.85, 1.42, accent="blue", title="Aprobacion", body="Estado, responsable, fechas y alertas de estancamiento.")
    add_card(slide, 5.09, 4.45, 1.85, 1.42, accent="amber", title="Cotizacion", body="Consolidado de opciones, valores y comparativos.")
    add_card(slide, 7.16, 4.45, 1.85, 1.42, accent="orange", title="Compra / entrega", body="Seguimiento de orden, recepcion y cumplimiento del proveedor.")
    add_card(slide, 9.23, 4.45, 2.25, 1.42, accent="teal", title="Uso final", body="Cierre con evidencia de destino, uso real y trazabilidad del recurso.")
    add_footer(slide, "Compras se presenta aqui como un proceso corporativo completo: solicitud, aprobacion, cotizacion, compra, entrega y uso final.")

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Aportes y mejoras en compras", "Mas control, menos tiempos muertos y mejor lectura del estado real del proceso")
    add_picture(slide, asset_paths["purchase_chart"], 7.02, 1.58, width=5.45)
    add_bullets(
        slide,
        0.95,
        1.72,
        5.65,
        4.45,
        [
            "Trazabilidad completa desde que nace la solicitud hasta que el recurso llega y se usa.",
            "Alertas cuando una aprobacion se estanca, una cotizacion no llega o una entrega se retrasa.",
            "Comparacion mas rapida de proveedores, valores, tiempos y condiciones.",
            "Seguimiento de compra contra recepcion y uso final, no solo hasta la orden.",
            "Mejor control para auditoria, presupuesto, cumplimiento y servicio al area solicitante.",
        ],
        font_size=15,
    )
    add_stat_card(slide, 0.95, 6.05, 1.75, "Control", "End to end", "violet")
    add_stat_card(slide, 2.9, 6.05, 1.75, "Alertas", "Tempranas", "orange")
    add_stat_card(slide, 4.85, 6.05, 1.75, "Cierre", "Con evidencia", "teal")
    add_footer(slide, "Grafico referencial: muestra como una gestion mas inteligente puede elevar la visibilidad y el control del proceso de compras.")

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "soft")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Valor para Operaciones y Productividad", "La mayor ganancia aparece cuando las areas se leen juntas")
    add_bullets(
        slide,
        0.95,
        1.8,
        5.45,
        3.65,
        [
            "Una cuadrilla puede verse afectada por ausentismo, falta de material, demoras de compra o indisponibilidad de vehiculo.",
            "El jefe de operaciones necesita entender esas causas juntas y no en reportes separados.",
            "El sistema puede priorizar frentes en riesgo, cuellos de botella y acciones correctivas con mas rapidez.",
        ],
        font_size=16,
    )

    nodes = [
        ("Personal", 7.0, 1.95, "violet"),
        ("Materiales", 9.2, 1.95, "amber"),
        ("Compras", 11.0, 3.1, "blue"),
        ("Transporte", 9.2, 4.2, "orange"),
        ("Operacion", 7.0, 4.2, "teal"),
    ]
    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.55), Inches(2.76), Inches(1.8), Inches(1.2))
    center.fill.solid()
    center.fill.fore_color.rgb = rgb("midnight")
    center.line.fill.background()
    add_textbox(slide, 8.85, 3.08, 1.18, 0.2, "Decision", font_size=16, bold=True, color="white", align=PP_ALIGN.CENTER)

    for label, lx, ly, accent_name in nodes:
        bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(ly), Inches(1.55), Inches(0.82))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = rgb("white")
        bubble.line.color.rgb = rgb(accent_name)
        bubble.line.width = Pt(2.0)
        add_textbox(slide, lx + 0.14, ly + 0.23, 1.25, 0.18, label, font_size=13, bold=True, color="midnight", align=PP_ALIGN.CENTER)
        add_arrow(slide, lx + 0.78, ly + 0.4, 9.45, 3.35, accent_name)

    add_textbox(slide, 7.0, 5.55, 5.25, 0.48, "Cuando la empresa conecta personas, soportes, materiales, compras y transporte, la productividad deja de verse por intuicion y empieza a verse con datos.", font_size=15, color="midnight")
    add_footer(slide, "El valor corporativo mas alto aparece en los cruces entre procesos, no solo dentro de cada area aislada.")

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "white")
    add_brand_marks(slide, asset_paths)
    add_title(slide, "Valor para Gerencia, Finanzas, Juridica e Indicadores", "Una sola capa de lectura ejecutiva sobre multiples procesos")
    add_stat_card(slide, 0.94, 1.62, 2.25, "Costos y desvio", "Visibles", "orange")
    add_stat_card(slide, 3.35, 1.62, 2.25, "Cumplimiento", "Trazable", "blue")
    add_stat_card(slide, 5.76, 1.62, 2.25, "Riesgos", "Priorizables", "teal")
    add_stat_card(slide, 8.17, 1.62, 2.25, "Indicadores", "Conversacionales", "violet")
    add_stat_card(slide, 10.58, 1.62, 1.84, "Soportes", "Leibles", "magenta")
    add_bullets(
        slide,
        0.96,
        3.0,
        5.8,
        2.5,
        [
            "Finanzas puede leer gastos, compras, consumos y desviaciones con mayor contexto operativo.",
            "Juridica puede resumir documentos, fechas, obligaciones y trazabilidad de casos o compromisos.",
            "Gerencia puede pedir KPI, alertas, causas probables y priorizacion sin esperar consolidaciones manuales.",
        ],
        font_size=15,
    )

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(3.0), Inches(5.0), Inches(2.35))
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb("soft")
    panel.line.color.rgb = rgb("line")
    panel.line.width = Pt(1.0)
    add_textbox(slide, 7.18, 3.18, 3.4, 0.24, "Preguntas ejecutivas que habilita", font_size=16, bold=True, color="midnight")
    add_bullets(
        slide,
        7.18,
        3.58,
        4.4,
        1.35,
        [
            "Que frentes tienen mayor riesgo esta semana?",
            "Donde se concentra el costo o el atraso?",
            "Que pendientes deben intervenir primero los jefes?",
        ],
        font_size=13,
    )
    add_footer(slide, "La vision de largo plazo es una empresa que no solo guarda datos, sino que los entiende y los usa para decidir.")

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, "midnight")
    add_brand_marks(slide, asset_paths, dark=True)
    add_title(slide, "Ruta de crecimiento", "De base analitica actual a copiloto corporativo para empresa grande", light=True)

    roadmap = [
        ("Fase 1", "Consolidar lo actual", "Empleados y ausentismo con uso gerencial y visibilidad real.", "violet"),
        ("Fase 2", "Expandir procesos", "Logistica, transporte, compras y tratamiento documental.", "blue"),
        ("Fase 3", "Empresa conectada", "Cruces operativos, alertas, seguimiento y lectura ejecutiva integral.", "teal"),
    ]
    x = 0.92
    for idx, (phase, title, body, accent_name) in enumerate(roadmap):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(3.55), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("white")
        card.line.fill.background()
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(2.28), Inches(1.0), Inches(0.28))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(accent_name)
        band.line.fill.background()
        add_textbox(slide, x + 0.34, 2.34, 0.7, 0.1, phase, font_size=9, bold=True, color="white")
        add_textbox(slide, x + 0.18, 2.8, 3.0, 0.28, title, font_size=18, bold=True, color="midnight")
        add_textbox(slide, x + 0.18, 3.24, 3.0, 0.7, body, font_size=12, color="slate")
        if idx < len(roadmap) - 1:
            add_arrow(slide, x + 3.55, 3.31, x + 3.95, 3.31, "line")
        x += 4.1

    add_textbox(slide, 1.0, 5.35, 11.2, 0.42, "Mensaje final: no se trata de tener otro chat, sino de darle a la empresa una nueva forma de leer, controlar y dirigir su operacion.", font_size=17, bold=True, color="white", align=PP_ALIGN.CENTER)
    add_picture(slide, asset_paths["ia_icon"], 5.78, 5.92, height=0.9)
    add_footer(slide, "Presentacion v2 generada con branding y graficos referenciales para conversacion directiva.")

    prs.save(PPTX_PATH)


def main() -> None:
    ensure_dirs()
    logos = prepare_brand_assets()
    domain_rows = load_domain_statuses()
    asset_paths = {
        **logos,
        "cover_bg": create_cover_background(),
        "domain_chart": create_domain_maturity_chart(domain_rows),
        "absence_chart": create_absence_chart(),
        "inventory_chart": create_inventory_chart(),
        "transport_chart": create_transport_chart(),
        "document_chart": create_document_chart(),
        "purchase_chart": create_purchase_impact_chart(),
    }
    build_presentation(asset_paths=asset_paths, domain_rows=domain_rows)
    print(f"PPTX generado en: {PPTX_PATH}")


if __name__ == "__main__":
    main()
